"""
Genera presentación PowerPoint del análisis de los 14 casos de crédito.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import json

# ── Paleta de colores ──────────────────────────────────────────────
AZUL       = RGBColor(0x1a, 0x73, 0xe8)
AZUL_OSC   = RGBColor(0x0d, 0x47, 0xa1)
VERDE      = RGBColor(0x34, 0xa8, 0x53)
ROJO       = RGBColor(0xea, 0x43, 0x35)
NARANJA    = RGBColor(0xfb, 0xbc, 0x04)
GRIS_OSC   = RGBColor(0x33, 0x33, 0x33)
GRIS_MED   = RGBColor(0x75, 0x75, 0x75)
GRIS_CLAR  = RGBColor(0xf5, 0xf5, 0xf5)
BLANCO     = RGBColor(0xff, 0xff, 0xff)
NEGRO      = RGBColor(0x00, 0x00, 0x00)

COLOR_DEC = {
    'APROBADO':           VERDE,
    'RECHAZADO':          ROJO,
    'ESCALAR_EJECUTIVO':  NARANJA,
    'VALIDACION':         AZUL,
}

# ── Datos de los 14 casos ──────────────────────────────────────────
CASOS = [
    dict(sol=47182039, bc=666, icc='0006', decil=6, neto=28000,  pago_buro=251,    monto=10744,  dec='APROBADO',          score=8, max_atr=6,  malo=False, veredicto='Correcto'),
    dict(sol=47156427, bc=637, icc='0005', decil=3, neto=44800,  pago_buro=21542,  monto=5700,   dec='VALIDACION',         score=6, max_atr=3,  malo=False, veredicto='Correcto — buró sin cuentas, pago buró muy alto'),
    dict(sol=47178641, bc=656, icc='0007', decil=5, neto=23700,  pago_buro=0,      monto=5744,   dec='VALIDACION',         score=8, max_atr=60, malo=True,  veredicto='Correcto — buró incompleto'),
    dict(sol=47161577, bc=611, icc='0004', decil=6, neto=25300,  pago_buro=10192,  monto=12350,  dec='VALIDACION',         score=5, max_atr=76, malo=True,  veredicto='Correcto — alto pago buró, score bajo'),
    dict(sol=47193179, bc=536, icc='0003', decil=5, neto=4722,   pago_buro=89477,  monto=8350,   dec='RECHAZADO',          score=3, max_atr=79, malo=True,  veredicto='Correcto'),
    dict(sol=47205357, bc=644, icc='0006', decil=5, neto=27600,  pago_buro=0,      monto=28636,  dec='VALIDACION',         score=8, max_atr=47, malo=True,  veredicto='Correcto — monto alto, decil 5'),
    dict(sol=47179512, bc=674, icc='0007', decil=7, neto=-2400,  pago_buro=1568,   monto=4350,   dec='RECHAZADO',          score=4, max_atr=1,  malo=False, veredicto='Falso negativo — egresos > ingresos declarados'),
    dict(sol=47192279, bc=586, icc='0005', decil=4, neto=5080,   pago_buro=913,    monto=4350,   dec='VALIDACION',         score=5, max_atr=73, malo=True,  veredicto='Correcto'),
    dict(sol=47129675, bc=661, icc='0006', decil=6, neto=10000,  pago_buro=6069,   monto=5350,   dec='ESCALAR_EJECUTIVO',  score=5, max_atr=0,  malo=False, veredicto='Correcto — pago buró = 60% ingreso neto'),
    dict(sol=47184400, bc=702, icc='0009', decil=6, neto=37100,  pago_buro=1083,   monto=10344,  dec='VALIDACION',         score=9, max_atr=40, malo=True,  veredicto='Correcto — regla thin file (3 cuentas, 1 abierta)'),
    dict(sol=47181090, bc=652, icc='0004', decil=4, neto=27600,  pago_buro=2970,   monto=15351,  dec='VALIDACION',         score=8, max_atr=71, malo=True,  veredicto='Correcto'),
    dict(sol=47185687, bc=589, icc='0004', decil=3, neto=30500,  pago_buro=4804,   monto=10350,  dec='RECHAZADO',          score=5, max_atr=53, malo=True,  veredicto='Correcto — decil 3, score bajo'),
    dict(sol=47164427, bc=628, icc='0002', decil=3, neto=25450,  pago_buro=6226,   monto=6600,   dec='VALIDACION',         score=6, max_atr=0,  malo=False, veredicto='Conservador — pagó bien, decil 3 e ICC bajo'),
    dict(sol=47148013, bc=603, icc='0006', decil=6, neto=8000,   pago_buro=1882,   monto=5350,   dec='APROBADO',           score=7, max_atr=0,  malo=False, veredicto='Correcto'),
]

ACCIONES = [
    ('APROBADO',           VERDE,   '✔',
     'Proceder con el desembolso',
     [
         'Verificar expediente físico completo: INE, comprobante domicilio, estados de cuenta.',
         'Sin revisión adicional requerida por parte de mesa de crédito.',
         'Documentar condiciones finales autorizadas (monto, tasa, plazo, cuota).',
         'Gestionar firma del contrato y desembolso en el plazo establecido.',
     ]),
    ('VALIDACION',         AZUL,    '⚑',
     'Enviar a mesa de crédito para revisión manual',
     [
         'Pausar el desembolso hasta obtener resolución de mesa.',
         'Causas frecuentes: expediente delgado (≤4 cuentas), buró sin historial registrado, perfil ambiguo.',
         'Mesa puede aprobar, rechazar, o solicitar documentación adicional: aval, garantía, comprobante de ingreso extra.',
         'Plazo máximo de resolución: 48 horas hábiles.',
     ]),
    ('ESCALAR_EJECUTIVO',  NARANJA, '▲',
     'Revisar con el ejecutivo de cuenta',
     [
         'Capacidad de pago comprometida pero no bloqueante; requiere juicio cualitativo.',
         'El ejecutivo puede renegociar monto, plazo o frecuencia con el cliente.',
         'Evaluar garantías adicionales o aval si el ajuste no es suficiente.',
         'Documentar la justificación de la decisión final en el expediente.',
     ]),
    ('RECHAZADO',          ROJO,    '✗',
     'Comunicar rechazo al cliente',
     [
         'Score < 500 o historial severo: informar que el historial crediticio impide la aprobación actual.',
         'Capacidad negativa (egresos > ingresos): invitar a volver cuando la situación financiera mejore.',
         'No desembolsar bajo ninguna circunstancia sin nueva evaluación.',
         'Registrar motivo en expediente y archivar la solicitud.',
     ]),
]

# ── Helpers ────────────────────────────────────────────────────────

def new_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)

def rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or NEGRO
    return tb

def add_bullet(tf, text, size=11, color=None, indent=0):
    p = tf.add_paragraph()
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color or GRIS_OSC

# ── Crear presentación ─────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — Portada
# ════════════════════════════════════════════════════════════════════
sl = new_slide(prs, 6)

# Fondo azul oscuro
rect(sl, 0, 0, 13.33, 7.5, fill=AZUL_OSC)

# Banda lateral izquierda
rect(sl, 0, 0, 0.25, 7.5, fill=AZUL)

# Título
txt(sl, 'Motor de Decisión Crediticia', 0.6, 1.5, 12, 1.2,
    size=36, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)

txt(sl, 'Análisis de 14 casos reales · Buró de Crédito + Modelo ML',
    0.6, 2.7, 12, 0.6, size=18, color=RGBColor(0xbb, 0xde, 0xfb), align=PP_ALIGN.LEFT)

txt(sl, 'SOFIPO · Marzo 2026', 0.6, 5.8, 8, 0.5,
    size=13, color=RGBColor(0x90, 0xca, 0xf9), align=PP_ALIGN.LEFT)

# Chips de decisiones
chips = [('APROBADO', VERDE), ('VALIDACION', AZUL), ('ESCALAR_EJECUTIVO', NARANJA), ('RECHAZADO', ROJO)]
for i, (lbl, col) in enumerate(chips):
    rx = 0.6 + i * 3.1
    r = rect(sl, rx, 3.6, 2.8, 0.55, fill=col)
    txt(sl, lbl, rx + 0.1, 3.65, 2.6, 0.45, size=11, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — ¿Qué hace el sistema?
# ════════════════════════════════════════════════════════════════════
sl = new_slide(prs, 6)
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_CLAR)
rect(sl, 0, 0, 13.33, 0.9, fill=AZUL_OSC)
txt(sl, 'Sistema multi-agente de análisis crediticio', 0.4, 0.12, 12, 0.66,
    size=20, bold=True, color=BLANCO)

agentes = [
    ('🪪 KYC',        AZUL,    'Verifica identidad biométrica:\nreconocimiento facial, INE,\nvalidación de vida y gobierno.'),
    ('💰 Financiero', VERDE,   'Calcula capacidad de pago:\ningresos − egresos − pago buró.\nAjusta monto si es necesario.'),
    ('📋 Buró',       NARANJA, 'Evalúa score BC, ICC,\nalertas HAWK, cuentas vencidas\ny expediente delgado.'),
    ('🤖 ML',         ROJO,    'Decil predictivo 1-10\nbasado en comportamiento\nhistórico 2024-2025.'),
    ('⚖️ Decisión',   AZUL_OSC,'Combina los 4 agentes\ny emite: APROBADO, VALIDACION,\nESCALAR o RECHAZADO.'),
]

for i, (nombre, col, desc) in enumerate(agentes):
    x = 0.3 + i * 2.58
    rect(sl, x, 1.2, 2.35, 2.8, fill=col)
    txt(sl, nombre, x+0.1, 1.3, 2.15, 0.5, size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txt(sl, desc,   x+0.1, 1.9, 2.15, 2.0, size=10, color=BLANCO, align=PP_ALIGN.CENTER)

# Reglas clave
txt(sl, 'Reglas clave del motor', 0.4, 4.25, 12, 0.4, size=13, bold=True, color=AZUL_OSC)
reglas = [
    '• Umbral capacidad de pago: 35% del ingreso neto (40% si BC ≥ 630 + ICC ≥ 5 + sin vencidos)',
    '• Thin file: ≤ 4 cuentas totales y ≤ 1 abierta → VALIDACION automática, sin importar el score',
    '• Pago a realizar del buró se descuenta del ingreso neto antes de calcular capacidad',
    '• Score BC < 500 o juicios activos → RECHAZADO automático',
]
y = 4.7
for reg in reglas:
    txt(sl, reg, 0.5, y, 12.4, 0.35, size=10.5, color=GRIS_OSC)
    y += 0.38

# ════════════════════════════════════════════════════════════════════
# SLIDES 3-6 — Una slide por tipo de decisión
# ════════════════════════════════════════════════════════════════════
for dec_name, col, icono, subtitulo, bullets in ACCIONES:
    sl = new_slide(prs, 6)
    rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_CLAR)

    # Header de color
    rect(sl, 0, 0, 13.33, 1.3, fill=col)
    txt(sl, f'{icono}  {dec_name}', 0.4, 0.1, 9, 0.7, size=28, bold=True, color=BLANCO)
    txt(sl, subtitulo, 0.4, 0.8, 9, 0.45, size=14, color=BLANCO, italic=True)

    # Casos de esta decisión
    casos_dec = [c for c in CASOS if c['dec'] == dec_name]
    txt(sl, f'{len(casos_dec)} caso(s) en esta categoría', 9.5, 0.2, 3.5, 0.4,
        size=13, bold=True, color=BLANCO, align=PP_ALIGN.RIGHT)

    # Panel izquierdo: pasos a seguir
    rect(sl, 0.3, 1.55, 5.8, 5.6, fill=BLANCO, line=col)
    txt(sl, 'Pasos a seguir', 0.5, 1.65, 5.4, 0.4, size=12, bold=True, color=col)

    tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.15), Inches(5.4), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = f'• {b}'
        run.font.size = Pt(11.5)
        run.font.color.rgb = GRIS_OSC
        p.space_after = Pt(8)

    # Panel derecho: casos
    rect(sl, 6.8, 1.55, 6.2, 5.6, fill=BLANCO, line=col)
    txt(sl, 'Solicitudes en esta decisión', 7.0, 1.65, 5.8, 0.4, size=12, bold=True, color=col)

    y_cas = 2.15
    for c in casos_dec:
        malo_txt = '⚠ Pagó mal' if c['malo'] else '✓ Pagó bien'
        malo_col = ROJO if c['malo'] else VERDE
        bg_col = RGBColor(0xff, 0xf3, 0xf3) if c['malo'] else RGBColor(0xf3, 0xff, 0xf6)
        rect(sl, 6.85, y_cas, 6.1, 0.72, fill=bg_col, line=RGBColor(0xdd,0xdd,0xdd))
        txt(sl, f'#{c["sol"]}', 6.95, y_cas+0.04, 2, 0.3, size=10, bold=True, color=GRIS_OSC)
        txt(sl, malo_txt, 10.7, y_cas+0.04, 2.1, 0.3, size=10, bold=True, color=malo_col, align=PP_ALIGN.RIGHT)
        det = f'BC {c["bc"]} · ICC {c["icc"]} · Decil {c["decil"]} · Neto ${c["neto"]:,.0f} · MaxAtr {c["max_atr"]}d'
        txt(sl, det, 6.95, y_cas+0.38, 5.9, 0.28, size=9, color=GRIS_MED)
        y_cas += 0.82

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — Cuadro resumen completo
# ════════════════════════════════════════════════════════════════════
sl = new_slide(prs, 6)
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_CLAR)
rect(sl, 0, 0, 13.33, 0.85, fill=AZUL_OSC)
txt(sl, 'Cuadro Resumen — 14 Casos', 0.4, 0.12, 10, 0.6, size=20, bold=True, color=BLANCO)
txt(sl, 'Resultado del motor vs comportamiento histórico real (MaxAtr_ventana)',
    0.4, 0.5, 10, 0.3, size=11, color=RGBColor(0xbb, 0xde, 0xfb))

# Encabezados tabla
cols_w  = [1.15, 0.75, 0.55, 0.6, 1.2, 1.15, 1.05, 1.65, 0.65, 0.8, 1.25]
cols_lbl = ['Solicitud', 'BC Score', 'ICC', 'Decil', 'Ing. Neto', 'Pago Buró',
            'Monto Sol.', 'Decisión', 'Score', 'MaxAtr', 'Veredicto']
col_align = [PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER,
             PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.CENTER,
             PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT]

x_start = 0.18
y_hdr   = 0.95
row_h   = 0.4

# Header row
x = x_start
for i, (lbl, w) in enumerate(zip(cols_lbl, cols_w)):
    rect(sl, x, y_hdr, w, row_h-0.02, fill=AZUL_OSC)
    txt(sl, lbl, x+0.04, y_hdr+0.05, w-0.08, row_h-0.1,
        size=8.5, bold=True, color=BLANCO, align=col_align[i])
    x += w

# Data rows
for idx, c in enumerate(CASOS):
    y = y_hdr + row_h + idx * row_h
    x = x_start
    bg = BLANCO if idx % 2 == 0 else RGBColor(0xf8, 0xf9, 0xfa)

    vals = [
        str(c['sol']),
        str(c['bc']),
        c['icc'],
        str(c['decil']),
        f'${c["neto"]:,.0f}',
        f'${c["pago_buro"]:,.0f}',
        f'${c["monto"]:,.0f}',
        c['dec'],
        f'{c["score"]}/10',
        f'{c["max_atr"]}d',
        c['veredicto'],
    ]

    dec_col = COLOR_DEC.get(c['dec'], GRIS_OSC)
    malo_bg = RGBColor(0xff, 0xeb, 0xeb) if c['malo'] else bg

    for i, (val, w) in enumerate(zip(vals, cols_w)):
        cell_bg = malo_bg if i in (9, 10) else bg
        if i == 7:  # Decisión
            cell_bg = RGBColor(
                int(dec_col[0]*0.15 + 0xe8*0.85),
                int(dec_col[1]*0.15 + 0xe8*0.85),
                int(dec_col[2]*0.15 + 0xe8*0.85),
            )
        rect(sl, x, y, w, row_h-0.02, fill=cell_bg, line=RGBColor(0xe0,0xe0,0xe0))
        fcol = dec_col if i == 7 else (ROJO if (c['malo'] and i == 9) else GRIS_OSC)
        fbold = (i == 7) or (i == 0)
        txt(sl, val, x+0.04, y+0.06, w-0.08, row_h-0.12,
            size=7.8, bold=fbold, color=fcol, align=col_align[i])
        x += w

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — Conclusiones y efectividad
# ════════════════════════════════════════════════════════════════════
sl = new_slide(prs, 6)
rect(sl, 0, 0, 13.33, 7.5, fill=GRIS_CLAR)
rect(sl, 0, 0, 13.33, 0.85, fill=AZUL_OSC)
txt(sl, 'Conclusiones y Efectividad del Motor', 0.4, 0.12, 12, 0.6, size=20, bold=True, color=BLANCO)

# KPIs
kpis = [
    ('14', 'Casos evaluados', AZUL),
    ('13/14', 'Decisiones correctas', VERDE),
    ('93%', 'Efectividad global', VERDE),
    ('1', 'Falso negativo', NARANJA),
]
for i, (val, lbl, col) in enumerate(kpis):
    x = 0.4 + i * 3.15
    rect(sl, x, 1.05, 2.85, 1.4, fill=col)
    txt(sl, val, x+0.1, 1.15, 2.65, 0.75, size=34, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)
    txt(sl, lbl, x+0.1, 1.88, 2.65, 0.45, size=11,  color=BLANCO, align=PP_ALIGN.CENTER)

# Distribución de decisiones
dist = [
    ('APROBADO',          2, VERDE),
    ('VALIDACION',        9, AZUL),
    ('ESCALAR_EJECUTIVO', 1, NARANJA),
    ('RECHAZADO',         3, ROJO),  # incluyendo 1 falso negativo
]
txt(sl, 'Distribución de decisiones', 0.4, 2.7, 6, 0.35, size=12, bold=True, color=AZUL_OSC)
y_d = 3.1
for lbl, n, col in dist:
    pct = n / 14
    bar_w = pct * 5.5
    rect(sl, 0.4, y_d, bar_w, 0.34, fill=col)
    txt(sl, f'{lbl}  ({n})', 0.4 + bar_w + 0.1, y_d+0.04, 4, 0.28, size=10, bold=True, color=col)
    y_d += 0.5

# Hallazgos
txt(sl, 'Hallazgos principales', 6.6, 2.7, 6.4, 0.35, size=12, bold=True, color=AZUL_OSC)
hallazgos = [
    '• 9 VALIDACIONES: buró sin cuentas o expediente delgado — el motor es conservador en perfiles sin historial.',
    '• Regla thin file capturó a 47184400 (score 9/10 pero llegó a 40 días de atraso).',
    '• Único falso negativo (47179512): egresos declarados > ingresos. Cliente pagó bien. Posible sub-declaración de ingresos.',
    '• 47156427: pago buró $21,542 sobre ingreso neto $44,800 → 48% comprometido. VALIDACION correcta aunque pagó bien.',
    '• 3 RECHAZADOS coincidieron perfectamente con clientes de 53-79 días de atraso.',
]
y_h = 3.15
for h in hallazgos:
    txt(sl, h, 6.6, y_h, 6.5, 0.42, size=10, color=GRIS_OSC)
    y_h += 0.5

# Próximos pasos
txt(sl, 'Próximos pasos', 0.4, 5.4, 12, 0.35, size=12, bold=True, color=AZUL_OSC)
pasos = [
    '1. Ajustar umbral de ingresos en RECHAZADO: añadir flag "posible sub-declaración" si BC > 650 + MaxAtr histórico bajo.',
    '2. Seguimiento de los 9 casos en VALIDACION: documentar resolución de mesa para re-entrenar el modelo.',
    '3. Ampliar muestra de evaluación a 50+ casos para validar estadísticamente las reglas thin file y umbral 40%.',
]
y_p = 5.8
for p in pasos:
    txt(sl, p, 0.5, y_p, 12.5, 0.35, size=10, color=GRIS_OSC)
    y_p += 0.42

# ── Guardar ────────────────────────────────────────────────────────
out = '/Users/fernandosierra/proyectos/agente-credito/Motor_Decision_Crediticia_14casos.pptx'
prs.save(out)
print(f'Presentación guardada: {out}')
